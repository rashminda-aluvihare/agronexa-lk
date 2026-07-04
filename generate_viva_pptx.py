import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Base path for project diagrams
    base_dir = os.getcwd()
    diag_dir = os.path.join(base_dir, "diagrams")
    logo_path = os.path.join(base_dir, "logo.png")

    # Helper function to safely insert picture
    def safe_add_picture(slide, img_path, left, top, width=None, height=None):
        if os.path.exists(img_path):
            if width and height:
                return slide.shapes.add_picture(img_path, left, top, width=width, height=height)
            elif width:
                return slide.shapes.add_picture(img_path, left, top, width=width)
            elif height:
                return slide.shapes.add_picture(img_path, left, top, height=height)
            else:
                return slide.shapes.add_picture(img_path, left, top)
        else:
            print(f"Warning: Image not found at {img_path}")
            return None

    # Color Palette Constants
    COLOR_PRIMARY_BG   = RGBColor(15, 23, 42)    # Dark Navy Slate #0F172A
    COLOR_HEADER_BG    = RGBColor(20, 83, 45)    # Dark Green #14532D
    COLOR_ACCENT_GREEN = RGBColor(34, 197, 94) # Vibrant Green #22C55E
    COLOR_GOLD         = RGBColor(234, 179, 8)   # Amber Gold #EAB308
    COLOR_CARD_BG      = RGBColor(248, 250, 252) # Clean Light Slate #F8FAFC
    COLOR_CARD_BORDER  = RGBColor(226, 232, 240) # Slate border
    COLOR_TEXT_MAIN    = RGBColor(15, 23, 42)    # Dark slate text
    COLOR_TEXT_MUTED   = RGBColor(71, 85, 105)   # Muted gray text
    COLOR_WHITE        = RGBColor(255, 255, 255)

    def add_header(slide, title_text, category_text="SLIATE HNDIT FINAL PROJECT VIVA — AGRONEXA LK"):
        # Header banner
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
        banner.fill.solid()
        banner.fill.fore_color.rgb = COLOR_HEADER_BG
        banner.line.fill.background()

        tf = banner.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.8)
        tf.margin_top = Inches(0.15)
        
        # Category label
        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ACCENT_GREEN

        # Title label
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(21)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE

        # Add logo on header top right if exists
        safe_add_picture(slide, logo_path, Inches(11.8), Inches(0.18), height=Inches(0.8))

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.fill.background()
        return card

    # ==================== SLIDE 1: TITLE SLIDE ====================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_PRIMARY_BG
    bg1.line.fill.background()

    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(0.15), Inches(5.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT_GREEN
    bar.line.fill.background()

    txBox = slide1.shapes.add_textbox(Inches(1.2), Inches(1.0), Inches(8.2), Inches(5.5))
    tf1 = txBox.text_frame
    tf1.word_wrap = True

    p0 = tf1.paragraphs[0]
    p0.text = "SRI LANKA INSTITUTE OF ADVANCED TECHNOLOGICAL EDUCATION (SLIATE)"
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_GOLD

    p0_sub = tf1.add_paragraph()
    p0_sub.text = "Higher National Diploma in Information Technology (HNDIT) — ATI Dehiwala"
    p0_sub.font.size = Pt(11)
    p0_sub.font.color.rgb = COLOR_WHITE
    p0_sub.space_after = Pt(15)

    p1 = tf1.add_paragraph()
    p1.text = "AgroNexa LK"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE

    p2 = tf1.add_paragraph()
    p2.text = "Smart Farming Platform with Blockchain-Inspired Technology for Sri Lanka"
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_ACCENT_GREEN
    p2.space_after = Pt(20)

    p3 = tf1.add_paragraph()
    p3.text = "Final Report Defense Presentation | Submission Date: 30th June 2026"
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.space_after = Pt(25)

    p4 = tf1.add_paragraph()
    p4.text = "Student Name: W.W.M.R.H.W RASHMINDA ALUVIHARE\nReg. Number: DEH/IT/2324/F/0038\nSupervisor: Dr. DSK Mendis"
    p4.font.size = Pt(13)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_WHITE

    safe_add_picture(slide1, logo_path, Inches(9.6), Inches(1.8), width=Inches(3.2))

    # ==================== SLIDE 2: EXECUTIVE SUMMARY & BACKGROUND ====================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "1. Executive Summary & Project Background")

    c1 = add_card(slide2, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tf_c1 = c1.text_frame
    tf_c1.word_wrap = True
    tf_c1.margin_left = Inches(0.3)
    tf_c1.margin_top = Inches(0.3)

    p = tf_c1.paragraphs[0]
    p.text = "📖 Project Abstract & Context"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADER_BG

    abs_points = [
        ("Agricultural Importance: ", "Agriculture plays a vital role in Sri Lanka's economy, yet smallholder farmers lack direct market access and fair pricing."),
        ("Development Methodology: ", "Followed the Prototype Software Development Model for continuous user feedback and iterative improvements."),
        ("Core Capability: ", "Web-based B2B marketplace directly connecting farmers with corporate buyers, backed by PostgreSQL and Express.js."),
        ("Security Integration: ", "Implements OTP authentication, role-based access control, and a tamper-evident SHA-256 transaction ledger.")
    ]
    for h, b in abs_points:
        p = tf_c1.add_paragraph()
        p.space_before = Pt(10)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_TEXT_MAIN
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_TEXT_MUTED

    c2 = add_card(slide2, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tf_c2 = c2.text_frame
    tf_c2.word_wrap = True
    tf_c2.margin_left = Inches(0.3)
    tf_c2.margin_top = Inches(0.3)

    p = tf_c2.paragraphs[0]
    p.text = "💻 Development Environment Specs"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADER_BG

    env_points = [
        ("Hardware Setup: ", "Intel Core i3 12th Gen CPU, 16 GB RAM, 512 GB Storage."),
        ("Operating System: ", "Windows 11 Development Environment."),
        ("Backend Runtime: ", "Node.js v20.11.0 & Express.js REST API framework."),
        ("Database Engine: ", "PostgreSQL v15 (relational schema with client pooling pg.Pool)."),
        ("Version Control & IDE: ", "Git & Visual Studio Code (VS Code).")
    ]
    for h, b in env_points:
        p = tf_c2.add_paragraph()
        p.space_before = Pt(10)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_TEXT_MAIN
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 3: PROBLEM ANALYSIS ====================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "2. Problem Analysis: Existing System vs. AgroNexa LK")

    c3_1 = add_card(slide3, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tf3_1 = c3_1.text_frame
    tf3_1.word_wrap = True
    tf3_1.margin_left = Inches(0.3)
    tf3_1.margin_top = Inches(0.3)

    p = tf3_1.paragraphs[0]
    p.text = "❌ Existing System Limitations"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = RGBColor(185, 28, 28)

    old_items = [
        ("Intermediary Exploitation: ", "Farmers rely on middlemen cartels, earning <30% of final retail margins."),
        ("Post-Harvest Losses: ", "Poor buyer matching and delayed transport cause 30%-40% crop wastage."),
        ("Idle Farm Machinery: ", "Smallholders cannot afford equipment, while machinery owners face long idle periods."),
        ("Price Opacity: ", "Lack of real-time wholesale price feeds forces farmers to sell blindly.")
    ]
    for h, b in old_items:
        p = tf3_1.add_paragraph()
        p.space_before = Pt(10)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_TEXT_MAIN
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_TEXT_MUTED

    c3_2 = add_card(slide3, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tf3_2 = c3_2.text_frame
    tf3_2.word_wrap = True
    tf3_2.margin_left = Inches(0.3)
    tf3_2.margin_top = Inches(0.3)

    p = tf3_2.paragraphs[0]
    p.text = "✅ AgroNexa LK Proposed Solution"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADER_BG

    new_items = [
        ("Direct B2B Marketplace: ", "Eliminates middlemen by connecting verified farmers directly with corporate buyers."),
        ("Machinery Sharing Economy: ", "Rent tractors and harvesters with double-booking prevention concurrency controls."),
        ("Cryptographic Auditability: ", "SHA-256 blockchain-style chained transaction logging for lease contracts."),
        ("Transparent Price Feeds: ", "Automated HARTI daily crop price web scraper with simulation fallback.")
    ]
    for h, b in new_items:
        p = tf3_2.add_paragraph()
        p.space_before = Pt(10)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_TEXT_MAIN
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 4: REQUIREMENTS & SECURITY ====================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "3. System Requirements & Security Specifications")

    c4_1 = add_card(slide4, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tf4_1 = c4_1.text_frame
    tf4_1.word_wrap = True
    tf4_1.margin_left = Inches(0.3)
    tf4_1.margin_top = Inches(0.3)

    p = tf4_1.paragraphs[0]
    p.text = "👥 Stakeholder Analysis"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADER_BG

    sh_items = [
        ("Primary Stakeholders: ", "System Administrator, Local Farmer/Seller, B2B Corporate Buyer."),
        ("Secondary Stakeholders: ", "Agricultural Businesses, Government Authorities, Developers, University Supervisors (Dr. DSK Mendis)."),
        ("Client Devices: ", "Android 9+, iOS 13+, Windows, macOS via Chrome, Firefox, Safari, Edge."),
        ("Network Requirements: ", "Active 3G/4G/5G/Wi-Fi connection + Cellular SMS coverage for OTP alerts.")
    ]
    for h, b in sh_items:
        p = tf4_1.add_paragraph()
        p.space_before = Pt(10)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_TEXT_MAIN
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_TEXT_MUTED

    c4_2 = add_card(slide4, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tf4_2 = c4_2.text_frame
    tf4_2.word_wrap = True
    tf4_2.margin_left = Inches(0.3)
    tf4_2.margin_top = Inches(0.3)

    p = tf4_2.paragraphs[0]
    p.text = "🔒 Non-Functional & Security Specs"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADER_BG

    sec_items = [
        ("Password Encryption: ", "Bcrypt hashing algorithm with work factor 12."),
        ("Session Control: ", "JWT tokens signed with server secret, expiring after 24 hours."),
        ("Brute-Force Protection: ", "Locks account for 30 minutes after 5 consecutive failed login attempts."),
        ("API Rate Limiting: ", "Limits incoming requests to max 300 per 15 minutes per IP (express-rate-limit)."),
        ("DB Connection Pooling: ", "PostgreSQL client pool (pg.Pool) for optimized concurrent query execution.")
    ]
    for h, b in sec_items:
        p = tf4_2.add_paragraph()
        p.space_before = Pt(8)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_TEXT_MAIN
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 5: SYSTEM MODULES ====================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "4. Functional System Modules (9 Core Modules)")

    mods = [
        ("1. Authentication Module", "SMS OTP via Twilio with developer mock mode (123456), Bcrypt hashing, JWT tokens, 30-min lockout."),
        ("2. User KYC Verification", "Multer NIC front/back file upload pipeline with Admin Lightbox UI (Canvas 360° rotation & scaling)."),
        ("3. Crop Marketplace", "Publish, edit, deactivate stock listings; search and filter with interactive Leaflet.js map index."),
        ("4. Equipment Rental Module", "Machinery listing and reservation requests with SQL interval overlap validation."),
        ("5. Broadcast Request Network", "Corporate bulk buying requests triggering instant regional SMS broadcasts to district farmers."),
        ("6. Socket.IO Live Chat", "Real-time buyer-seller messaging with read receipts, double-tick indicators, and voice message player."),
        ("7. Cryptographic Ledger", "SHA-256 blockchain-style block chaining for confirmed rental agreements with <15ms audit traversal."),
        ("8. Daily Crop Price Index", "Axios & Cheerio HARTI web scraper with +/- 3% simulation fallback engine."),
        ("9. Tri-Lingual Engine", "Client-side translation dictionary (translations.js) supporting English, Sinhala, and Tamil seamlessly.")
    ]

    coords = [(0.8, 1.5), (4.8, 1.5), (8.8, 1.5),
              (0.8, 3.4), (4.8, 3.4), (8.8, 3.4),
              (0.8, 5.3), (4.8, 5.3), (8.8, 5.3)]

    for i, (title, desc) in enumerate(mods):
        left, top = coords[i]
        c = add_card(slide5, Inches(left), Inches(top), Inches(3.7), Inches(1.8))
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_HEADER_BG

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.space_before = Pt(4)
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 6: USE CASE DIAGRAM ====================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "5. System Scope & Use Case Architecture")

    c6_t = add_card(slide6, Inches(0.8), Inches(1.5), Inches(5.2), Inches(5.3))
    tf6 = c6_t.text_frame
    tf6.word_wrap = True
    tf6.margin_left = Inches(0.3)
    tf6.margin_top = Inches(0.3)

    p = tf6.paragraphs[0]
    p.text = "🎯 Use Case Scope Summary"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADER_BG

    uc_points = [
        ("Farmer Operations: ", "User registration, crop inventory listing, rental tool availability, broadcast response."),
        ("Buyer Operations: ", "Marketplace browse, Leaflet map filtering, equipment rental booking, district broadcast creation."),
        ("Admin Moderation: ", "KYC document review & approval/rejection, user account lockouts, ledger integrity checks, PDF report export."),
        ("System Boundary: ", "Strict role enforcement across REST API routes preventing unauthorized access.")
    ]
    for h, b in uc_points:
        p = tf6.add_paragraph()
        p.space_before = Pt(10)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_TEXT_MAIN
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_TEXT_MUTED

    add_card(slide6, Inches(6.3), Inches(1.5), Inches(6.2), Inches(5.3))
    safe_add_picture(slide6, os.path.join(diag_dir, "Figure_4.3_Use_Case_Diagram.png"), Inches(6.4), Inches(1.6), width=Inches(6.0))

    # ==================== SLIDE 7: DEPLOYMENT ARCHITECTURE ====================
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "6. Overall Architecture & Production Cloud Deployment")

    add_card(slide7, Inches(0.8), Inches(1.5), Inches(6.8), Inches(5.3))
    safe_add_picture(slide7, os.path.join(diag_dir, "Figure_4.1_Overall_System_Architecture.png"), Inches(0.9), Inches(1.6), width=Inches(6.6))

    c7 = add_card(slide7, Inches(7.8), Inches(1.5), Inches(4.7), Inches(5.3))
    tf7 = c7.text_frame
    tf7.word_wrap = True
    tf7.margin_left = Inches(0.3)
    tf7.margin_top = Inches(0.3)

    p = tf7.paragraphs[0]
    p.text = "☁️ Production Cloud Setup"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADER_BG

    dep_items = [
        ("Frontend Hosting: ", "Deployed on Vercel CDN Edge Networks for optimized global delivery."),
        ("Backend Hosting: ", "Hosted on Railway Node.js execution sandbox container node."),
        ("Database Hosting: ", "Railway Managed PostgreSQL instance (Port 5432)."),
        ("Cloud Storage: ", "Cloudinary CDN for secure NIC verification uploads and crop images."),
        ("Third-Party Gateways: ", "Twilio SMS API Sandbox and HARTI Scraper portal.")
    ]
    for h, b in dep_items:
        p = tf7.add_paragraph()
        p.space_before = Pt(10)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_TEXT_MAIN
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 8: DATABASE ERD ====================
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "7. Database Design & Entity Relationship Diagram (ERD)")

    c8_t = add_card(slide8, Inches(0.8), Inches(1.5), Inches(4.2), Inches(5.3))
    tf8 = c8_t.text_frame
    tf8.word_wrap = True
    tf8.margin_left = Inches(0.3)
    tf8.margin_top = Inches(0.3)

    p = tf8.paragraphs[0]
    p.text = "🗄️ Relational Schema"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADER_BG

    erd_points = [
        ("Core Tables: ", "users, crop_listings, equipment_bookings, rental_ledger, buyer_requests, direct_messages, audit_logs."),
        ("Data Integrity: ", "Foreign key constraints, non-null check constraints, compound date indexes."),
        ("Reputation Scoring: ", "Calculates seller score dynamically: min(5.0, (completed_rentals / 10.0) * 5.0).")
    ]
    for h, b in erd_points:
        p = tf8.add_paragraph()
        p.space_before = Pt(12)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_TEXT_MAIN
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_TEXT_MUTED

    add_card(slide8, Inches(5.2), Inches(1.5), Inches(7.3), Inches(5.3))
    safe_add_picture(slide8, os.path.join(diag_dir, "Figure_4.2_Entity_Relationship_Diagram.png"), Inches(5.3), Inches(1.6), width=Inches(7.1))

    # ==================== SLIDE 9: INNOVATION 1 - LEDGER ====================
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "8. Innovation I: Cryptographic SHA-256 Chained Ledger")

    c9_l = add_card(slide9, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tf9_l = c9_l.text_frame
    tf9_l.word_wrap = True
    tf9_l.margin_left = Inches(0.3)
    tf9_l.margin_top = Inches(0.3)

    p = tf9_l.paragraphs[0]
    p.text = "🔐 Blockchain-Inspired Ledger"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADER_BG

    led_points = [
        ("Concept: ", "Blockchain-inspired immutable transaction log for equipment lease contracts implemented in ledger.service.js."),
        ("Hash Chaining Formula: ", "block_hash = SHA256(tx_id + listing_id + renter_id + owner_id + amount + duration + agreement_hash + prev_hash)."),
        ("Audit Traversal Scan: ", "Sequential recalculation engine verifies entire chain in <15ms. Any manual DB row edit breaks the hash sequence immediately.")
    ]
    for h, b in led_points:
        p = tf9_l.add_paragraph()
        p.space_before = Pt(10)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_TEXT_MAIN
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_TEXT_MUTED

    add_card(slide9, Inches(6.7), Inches(1.5), Inches(5.8), Inches(5.3))
    safe_add_picture(slide9, os.path.join(diag_dir, "Figure_4.8_Cryptographic_Ledger_Workflow.png"), Inches(6.8), Inches(1.6), width=Inches(5.6))

    # ==================== SLIDE 10: INNOVATION 2 - CONCURRENCY ====================
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "9. Innovation II: Double-Booking Prevention Algorithm")

    add_card(slide10, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3))
    safe_add_picture(slide10, os.path.join(diag_dir, "Figure_4.5_Equipment_Booking_Sequence_Diagram.png"), Inches(0.9), Inches(1.6), width=Inches(5.6))

    c10_r = add_card(slide10, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), bg_color=COLOR_PRIMARY_BG, border_color=None)
    tf10_r = c10_r.text_frame
    tf10_r.word_wrap = True
    tf10_r.margin_left = Inches(0.3)
    tf10_r.margin_top = Inches(0.3)

    p = tf10_r.paragraphs[0]
    p.text = "⚡ Concurrency SQL Overlap Query"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD

    sql_code = [
        "SELECT id FROM equipment_bookings",
        "WHERE listing_id = $1 AND status = 'confirmed'",
        "  AND NOT (end_date < $2 OR start_date > $3);",
        "",
        "// Handling in equipment.controller.js:",
        "if (overlap.rows.length > 0) {",
        "  return res.status(409).json({",
        "    error: 'Selected dates overlap with an existing booking'",
        "  });",
        "}"
    ]
    for line in sql_code:
        p_sql = tf10_r.add_paragraph()
        p_sql.text = line
        p_sql.font.size = Pt(11)
        p_sql.font.name = "Consolas"
        p_sql.font.color.rgb = COLOR_ACCENT_GREEN if "SELECT" in line or "WHERE" in line else COLOR_WHITE

    # ==================== SLIDE 11: KYC & ADMIN MODERATION ====================
    slide11 = prs.slides.add_slide(blank_layout)
    add_header(slide11, "10. KYC Identity Verification & Admin Moderation")

    add_card(slide11, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3))
    safe_add_picture(slide11, os.path.join(diag_dir, "Figure_4.11_KYC_Verification_Activity_Diagram.png"), Inches(0.9), Inches(1.6), width=Inches(5.6))

    c11_r = add_card(slide11, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tf11_r = c11_r.text_frame
    tf11_r.word_wrap = True
    tf11_r.margin_left = Inches(0.3)
    tf11_r.margin_top = Inches(0.3)

    p = tf11_r.paragraphs[0]
    p.text = "🛡️ Admin Lightbox & Audit Tools"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADER_BG

    kyc_items = [
        ("Multer File Security: ", "Sanitizes upload filenames to prevent directory traversal and limits size to 5MB."),
        ("Canvas Lightbox Modal: ", "Interactive 360° image rotation and viewport scaling in admin portal for NIC verification."),
        ("Rejection Feedback: ", "Admins append rejection reasons (e.g. NIC blurry) sent via alert upon login attempt."),
        ("Audit Logs & PDF Exporter: ", "Captures admin actions and exports formatted PDF registers.")
    ]
    for h, b in kyc_items:
        p = tf11_r.add_paragraph()
        p.space_before = Pt(10)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_TEXT_MAIN
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 12: USER DASHBOARDS ====================
    slide12 = prs.slides.add_slide(blank_layout)
    add_header(slide12, "11. User Interface Dashboards & User Experience (UX)")

    add_card(slide12, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.3))
    safe_add_picture(slide12, os.path.join(diag_dir, "low_fidelity_seller_dashboard.png"), Inches(0.9), Inches(1.6), width=Inches(5.5))

    add_card(slide12, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    safe_add_picture(slide12, os.path.join(diag_dir, "low_fidelity_buyer_dashboard.png"), Inches(6.9), Inches(1.6), width=Inches(5.5))

    # ==================== SLIDE 13: TESTING & EMPIRICAL RESULTS ====================
    slide13 = prs.slides.add_slide(blank_layout)
    add_header(slide13, "12. Empirical Testing & Validation Results")

    test_cards = [
        ("🧪 Test Suite Results", "Executed 15 structured test cases across Authentication, KYC, Marketplace, Rentals, Cryptographic Ledger, and Chat.\n\n• Total Tests: 15\n• Passed: 15 (100% Pass Rate)\n• Failed: 0 (0%)", Inches(0.8)),
        ("⚡ System Latency Metrics", "• SHA-256 Block Mining: 12ms average.\n• Socket.IO Chat Latency: <45ms.\n• Ledger Traversal Scan: <15ms.\n• DB Query Execution: Accelerated via PostgreSQL indexing.", Inches(4.8)),
        ("📊 Lighthouse Audit Scores", "Google Lighthouse Performance Audit:\n\n• Mobile Performance: 91 / 100\n• Mobile SEO: 91 / 100\n• Desktop Performance: 97 / 100\n• Desktop SEO: 91 / 100", Inches(8.8))
    ]

    for title, desc, left_pos in test_cards:
        c = add_card(slide13, left_pos, Inches(1.5), Inches(3.7), Inches(5.3))
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_HEADER_BG

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.space_before = Pt(14)
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 14: FINANCIAL STATEMENT ====================
    slide14 = prs.slides.add_slide(blank_layout)
    add_header(slide14, "13. Project Financial Statement")

    c14 = add_card(slide14, Inches(2.2), Inches(1.8), Inches(8.933), Inches(4.7))
    tf14 = c14.text_frame
    tf14.word_wrap = True
    tf14.margin_left = Inches(0.5)
    tf14.margin_top = Inches(0.5)

    p = tf14.paragraphs[0]
    p.text = "💰 Project Expenditure Breakdown"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADER_BG

    p1 = tf14.add_paragraph()
    p1.text = "Backend & Database Cloud Hosting (Railway Sandbox):"
    p1.space_before = Pt(20)
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_MAIN

    p2 = tf14.add_paragraph()
    p2.text = "$5 USD per month  ×  4 Months  =  6,460 LKR"
    p2.space_before = Pt(10)
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_GOLD

    p3 = tf14.add_paragraph()
    p3.text = "Total Development & Hosting Expenditure: 6,460 LKR"
    p3.space_before = Pt(30)
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_HEADER_BG

    # ==================== SLIDE 15: LIMITATIONS ====================
    slide15 = prs.slides.add_slide(blank_layout)
    add_header(slide15, "14. System Limitations & Constraints")

    c15_1 = add_card(slide15, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tf15_1 = c15_1.text_frame
    tf15_1.word_wrap = True
    tf15_1.margin_left = Inches(0.3)
    tf15_1.margin_top = Inches(0.3)

    p = tf15_1.paragraphs[0]
    p.text = "⚠️ Architectural Limitations"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = RGBColor(185, 28, 28)

    lim_1 = [
        ("Centralized DB Dependency: ", "Ledger resides inside centralized PostgreSQL; direct DB console access could alter rows."),
        ("Manual KYC Bottleneck: ", "NIC review requires manual admin inspection, introducing human verification latency."),
        ("Internet-Only Client: ", "Requires continuous mobile data, restricting access for remote farmers with basic feature phones.")
    ]
    for h, b in lim_1:
        p = tf15_1.add_paragraph()
        p.space_before = Pt(10)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_TEXT_MAIN
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_TEXT_MUTED

    c15_2 = add_card(slide15, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tf15_2 = c15_2.text_frame
    tf15_2.word_wrap = True
    tf15_2.margin_left = Inches(0.3)
    tf15_2.margin_top = Inches(0.3)

    p = tf15_2.paragraphs[0]
    p.text = "⚠️ Operational Constraints"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = RGBColor(185, 28, 28)

    lim_2 = [
        ("Outbound-Only SMS Alerts: ", "SMS gateway is unidirectional; farmers receive text alerts but cannot reply offline."),
        ("Absence of Integrated Payment: ", "Relies on cash-on-delivery or manual bank transfers rather than live escrow IPG."),
        ("API Sandbox Sandbox Limits: ", "Twilio Developer Sandbox restricts SMS dispatches to pre-verified numbers.")
    ]
    for h, b in lim_2:
        p = tf15_2.add_paragraph()
        p.space_before = Pt(10)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_TEXT_MAIN
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 16: CONCLUSION & FUTURE WORK ====================
    slide16 = prs.slides.add_slide(blank_layout)
    add_header(slide16, "15. Conclusion & Future Work Roadmap")

    c16_1 = add_card(slide16, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tf16_1 = c16_1.text_frame
    tf16_1.word_wrap = True
    tf16_1.margin_left = Inches(0.3)
    tf16_1.margin_top = Inches(0.3)

    p = tf16_1.paragraphs[0]
    p.text = "🎉 Project Achievements Summary"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADER_BG

    ach_items = [
        "Engineered a scalable B2B marketplace bridging Sri Lanka's agrarian supply chain gaps.",
        "Eliminated middleman price exploitation through direct buyer-seller trading.",
        "Increased equipment utilization via an integrated machinery sharing economy.",
        "Guaranteed contract immutability using SHA-256 cryptographic chaining.",
        "Achieved 100% test pass rate and 92% System Usability Scale (SUS) score."
    ]
    for item in ach_items:
        p = tf16_1.add_paragraph()
        p.text = "✓  " + item
        p.space_before = Pt(10)
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN

    c16_2 = add_card(slide16, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tf16_2 = c16_2.text_frame
    tf16_2.word_wrap = True
    tf16_2.margin_left = Inches(0.3)
    tf16_2.margin_top = Inches(0.3)

    p = tf16_2.paragraphs[0]
    p.text = "🚀 Future Work Roadmap"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADER_BG

    fut_items = [
        ("Escrow Payment Gateway: ", "PayHere / DirectPay integration holding funds until machinery return."),
        ("Decentralized Ledger (DLT): ", "Hyperledger Fabric migration to eliminate single-point-of-failure."),
        ("Automated AI KYC: ", "OCR and Computer Vision models for automatic NIC parsing."),
        ("IPFS Document Storage: ", "Decentralized peer-to-peer storage for user identity files."),
        ("Interactive USSD & Offline SMS: ", "Allows feature phone users to trade without internet."),
        ("IoT Telemetry Sensors: ", "GPS and runtime logging IoT sensors on rental equipment.")
    ]
    for h, b in fut_items:
        p = tf16_2.add_paragraph()
        p.space_before = Pt(6)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = COLOR_TEXT_MAIN
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(11)
        r2.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 17: Q&A / THANK YOU ====================
    slide17 = prs.slides.add_slide(blank_layout)
    bg17 = slide17.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg17.fill.solid()
    bg17.fill.fore_color.rgb = COLOR_PRIMARY_BG
    bg17.line.fill.background()

    txBox17 = slide17.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.8))
    tf17 = txBox17.text_frame
    tf17.word_wrap = True

    p = tf17.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Thank You!"
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_GREEN

    p2 = tf17.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "Questions & Answers (Viva Defense)"
    p2.space_before = Pt(15)
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLOR_GOLD

    p3 = tf17.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.text = "AgroNexa LK — Smart Farming Platform with Blockchain-Inspired Technology"
    p3.space_before = Pt(25)
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(203, 213, 225)

    p4 = tf17.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    p4.text = "Student: W.W.M.R.H.W Rashminda Aluvihare (DEH/IT/2324/F/0038)\nSupervisor: Dr. DSK Mendis | SLIATE ATI Dehiwala"
    p4.space_before = Pt(20)
    p4.font.size = Pt(13)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_WHITE

    output_path = os.path.join(os.getcwd(), "AgroNexa_LK_Viva_Presentation.pptx")
    prs.save(output_path)
    print(f"Master presentation deck created successfully from Final Report at: {output_path}")

if __name__ == "__main__":
    create_deck()

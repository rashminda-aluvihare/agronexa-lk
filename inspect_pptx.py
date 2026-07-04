import sys
import os
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

def inspect():
    path = os.path.join(os.getcwd(), "AgroNexa_LK_Viva_Presentation.pptx")
    if not os.path.exists(path):
        print(f"File not found: {path}")
        return

    prs = Presentation(path)
    print(f"Total Slides in Deck: {len(prs.slides)}\n")

    for idx, slide in enumerate(prs.slides, 1):
        print(f"=== SLIDE {idx} ===")
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    txt = paragraph.text.strip()
                    if txt:
                        texts.append(txt)
        for t in texts:
            print(f"  - {t}")
        print()

if __name__ == "__main__":
    inspect()
